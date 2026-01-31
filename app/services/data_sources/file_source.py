from pathlib import Path
from typing import List, Dict, Any, Optional
import csv
import io
import logging

from .base import BaseDataSource

logger = logging.getLogger(__name__)


class FileSource(BaseDataSource):
    """Extract documents from local files with comprehensive format support"""

    def __init__(self):
        super().__init__("FileSource")
        self.supported_formats = {
            '.txt', '.md',           # Plain text
            '.pdf',                   # PDF (with OCR support)
            '.docx',                  # Word documents
            '.pptx',                  # PowerPoint
            '.xlsx', '.xls',          # Excel
            '.csv',                   # CSV
            '.html', '.htm',          # HTML
            '.rtf',                   # Rich Text Format
            '.png', '.jpg', '.jpeg', '.tiff', '.bmp',  # Images (OCR)
        }
        self._ocr_available = None
        self._markitdown_available = None

    def _check_markitdown_available(self) -> bool:
        """Check if markitdown is available for enhanced extraction"""
        if self._markitdown_available is None:
            try:
                from markitdown import MarkItDown
                self._markitdown_available = True
            except ImportError:
                self._markitdown_available = False
                logger.debug("markitdown not available, using built-in extractors")
        return self._markitdown_available

    def _extract_with_markitdown(self, path: Path) -> Optional[str]:
        """Try to extract content using markitdown library"""
        if not self._check_markitdown_available():
            return None
        try:
            from markitdown import MarkItDown
            md = MarkItDown()
            result = md.convert(str(path))
            if result and result.text_content:
                return result.text_content.strip()
        except Exception as e:
            logger.debug(f"markitdown extraction failed for {path}: {e}")
        return None

    def _check_ocr_available(self) -> bool:
        """Check if OCR dependencies are available"""
        if self._ocr_available is None:
            try:
                import pytesseract
                pytesseract.get_tesseract_version()
                self._ocr_available = True
            except Exception:
                self._ocr_available = False
                logger.warning("OCR not available: pytesseract or tesseract not installed")
        return self._ocr_available

    async def extract(self, **kwargs) -> List[Dict[str, Any]]:
        """Standard batch extraction (returns list)"""
        docs = []
        async for doc in self.extract_stream(**kwargs):
            docs.append(doc)
        return docs

    async def extract_stream(
        self,
        path: str = None,
        file_path: str = None,
        directory_path: str = None,
        ocr_enabled: bool = True,
        ocr_language: str = "eng",
        password: str = None,
        extract_tables: bool = True,
        extract_images: bool = False,
        **kwargs
    ):
        """
        Extract documents from files or directory as a stream
        """
        # Store extraction options
        self._ocr_enabled = ocr_enabled
        self._ocr_language = ocr_language
        self._password = password
        self._extract_tables = extract_tables
        self._extract_images = extract_images

        if path:
            # Auto-detect if file or directory
            p = Path(path)
            if p.is_file():
                file_path = path
            elif p.is_dir():
                directory_path = path

        if file_path:
            async for doc in self._extract_single_file_stream(file_path):
                yield doc
        elif directory_path:
            async for doc in self._extract_directory_stream(directory_path):
                yield doc
        else:
            raise ValueError("Must provide either 'path', 'file_path', or 'directory_path'")

    async def _extract_directory_stream(self, dir_path: str):
        """Scan directory and yield files one by one"""
        path = Path(dir_path)
        if not path.exists() or not path.is_dir():
            raise ValueError(f"Directory not found: {dir_path}")

        for file_path in path.rglob('*'):
            if file_path.suffix.lower() in self.supported_formats:
                try:
                    async for doc in self._extract_single_file_stream(str(file_path)):
                        yield doc
                except Exception as e:
                    logger.error(f"Error extracting {file_path}: {e}")
                    continue

    async def _extract_single_file_stream(self, file_path: str):
        """Stream content from a single file based on its extension"""
        path = Path(file_path).resolve()
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()
        if suffix not in self.supported_formats:
            return

        # Large files (PDF, CSV) are yielded page-by-page or batch-by-batch
        if suffix == '.pdf':
            async for doc in self._extract_pdf_stream(path):
                yield doc
        elif suffix == '.csv':
            async for doc in self._extract_csv_stream(path):
                yield doc
        else:
            # For other formats, fallback to standard extraction (yield single doc)
            docs = await self._extract_single_file(file_path)
            for doc in docs:
                yield doc

    async def _extract_directory(self, dir_path: str) -> List[Dict[str, Any]]:
        """Scan directory and extract all supported files"""
        path = Path(dir_path)

        if not path.exists() or not path.is_dir():
            raise ValueError(f"Directory not found: {dir_path}")

        documents = []

        for file_path in path.rglob('*'):
            if file_path.suffix.lower() in self.supported_formats:
                try:
                    docs = await self._extract_single_file(str(file_path))
                    documents.extend(docs)
                except Exception as e:
                    logger.error(f"Error extracting {file_path}: {e}")
                    continue

        logger.info(f"Extracted {len(documents)} documents from directory")
        return documents

    async def _extract_single_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Extract content from a single file"""
        path = Path(file_path).resolve()

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = path.suffix.lower()

        if suffix not in self.supported_formats:
            raise ValueError(f"Unsupported format: {suffix}")

        # Route to appropriate extractor
        content = ""
        extraction_method = "text"

        if suffix == '.pdf':
            content, extraction_method = self._extract_pdf(path)
        elif suffix == '.docx':
            content = self._extract_docx(path)
        elif suffix == '.pptx':
            content = self._extract_pptx(path)
        elif suffix in {'.xlsx', '.xls'}:
            content = self._extract_excel(path)
        elif suffix == '.csv':
            content = self._extract_csv(path)
        elif suffix in {'.html', '.htm'}:
            content = self._extract_html(path)
        elif suffix == '.rtf':
            content = self._extract_rtf(path)
        elif suffix in {'.png', '.jpg', '.jpeg', '.tiff', '.bmp'}:
            content = self._extract_image_ocr(path)
            extraction_method = "ocr"
        elif suffix in {'.txt', '.md'}:
            content = self._extract_text(path)
        else:
            return []

        if not content or not content.strip():
            logger.warning(f"No content extracted from: {file_path}")
            return []

        return [{
            'content': content,
            'metadata': {
                'source': self.source_name,
                'filename': path.name,
                'filepath': str(path.absolute()),
                'type': suffix[1:],
                'extraction_method': extraction_method
            }
        }]

    async def _extract_pdf_stream(self, path: Path):
        """Yield PDF content page-by-page"""
        try:
            import fitz
            doc = fitz.open(str(path))
            
            # Check for encryption
            if doc.is_encrypted:
                if self._password:
                    if not doc.authenticate(self._password):
                        raise ValueError("PDF password incorrect")
                else:
                    if not doc.authenticate(""):
                        raise ValueError("PDF is password protected")

            for page_num, page in enumerate(doc):
                page_text = page.get_text().strip()
                
                # Dynamic OCR fallback for large pages
                if len(page_text) < 50 and self._ocr_enabled and self._check_ocr_available():
                    page_text = self._ocr_pdf_page(page) or ""

                if page_text:
                    yield {
                        'content': page_text,
                        'metadata': {
                            'source': self.source_name,
                            'filename': path.name,
                            'page': page_num + 1,
                            'total_pages': len(doc),
                            'type': 'pdf'
                        }
                    }
            doc.close()
        except ImportError:
            # Fallback to batch if fitz missing (though ideally fitz should be there)
            content, method = self._extract_pdf(path)
            yield {
                'content': content,
                'metadata': {'source': self.source_name, 'filename': path.name, 'type': 'pdf', 'method': method}
            }

    async def _extract_csv_stream(self, path: Path, batch_size: int = 100):
        """Yield CSV content in batches of rows"""
        import csv
        for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
            try:
                with open(path, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    header = next(reader, None)
                    
                    rows = []
                    batch_idx = 0
                    for row_num, row in enumerate(reader, 1):
                        rows.append(row)
                        if len(rows) >= batch_size:
                            yield {
                                'content': self._format_table([header] + rows if header else rows),
                                'metadata': {
                                    'source': self.source_name,
                                    'filename': path.name,
                                    'batch': batch_idx,
                                    'type': 'csv'
                                }
                            }
                            rows = []
                            batch_idx += 1
                    
                    if rows:
                        yield {
                            'content': self._format_table([header] + rows if header else rows),
                            'metadata': {
                                'source': self.source_name,
                                'filename': path.name,
                                'batch': batch_idx,
                                'type': 'csv'
                            }
                        }
                return
            except (UnicodeDecodeError, StopIteration):
                continue

    def _extract_pdf(self, path: Path) -> tuple[str, str]:
        """
        Extract text from PDF with fallback to OCR for scanned documents

        Returns:
            Tuple of (content, extraction_method)
        """
        text_content = []
        extraction_method = "text"
        is_encrypted = False

        # First, try markitdown for better extraction (handles many edge cases)
        markitdown_content = self._extract_with_markitdown(path)
        if markitdown_content and len(markitdown_content) > 100:
            return markitdown_content, "markitdown"

        try:
            # Try pymupdf (fitz) first for better extraction
            import fitz

            try:
                doc = fitz.open(str(path))

                # Check if PDF is encrypted
                if doc.is_encrypted:
                    is_encrypted = True
                    if self._password:
                        if not doc.authenticate(self._password):
                            raise ValueError(f"PDF is password-protected. The provided password is incorrect.")
                    else:
                        # Try empty password (some PDFs have empty password protection)
                        if not doc.authenticate(""):
                            doc.close()
                            raise ValueError(f"PDF is password-protected. Please provide the 'password' parameter to extract content.")

                for page_num, page in enumerate(doc):
                    # Extract text
                    page_text = page.get_text()

                    # If page has very little text, might be scanned - try OCR
                    if len(page_text.strip()) < 50 and self._ocr_enabled and self._check_ocr_available():
                        ocr_text = self._ocr_pdf_page(page)
                        if ocr_text:
                            page_text = ocr_text
                            extraction_method = "ocr"

                    if page_text.strip():
                        text_content.append(page_text)

                    # Extract tables if enabled
                    if self._extract_tables:
                        try:
                            tables = page.find_tables()
                            for table in tables:
                                table_text = self._format_table(table.extract())
                                if table_text:
                                    text_content.append(f"\n[Table]\n{table_text}")
                        except Exception as te:
                            logger.debug(f"Table extraction failed for page {page_num}: {te}")

                doc.close()

            except fitz.fitz.FileDataError as e:
                if "encrypted" in str(e).lower() or "password" in str(e).lower():
                    raise ValueError(f"PDF is password-protected. Please provide the 'password' parameter to extract content.")
                raise

        except ImportError:
            # Fallback to PyPDF2
            logger.debug("pymupdf not available, using PyPDF2")
            import PyPDF2

            with open(path, 'rb') as file:
                try:
                    pdf_reader = PyPDF2.PdfReader(file)

                    # Check if encrypted
                    if pdf_reader.is_encrypted:
                        is_encrypted = True
                        if self._password:
                            if not pdf_reader.decrypt(self._password):
                                raise ValueError(f"PDF is password-protected. The provided password is incorrect.")
                        else:
                            # Try empty password
                            if not pdf_reader.decrypt(""):
                                raise ValueError(f"PDF is password-protected. Please provide the 'password' parameter to extract content.")

                except Exception as e:
                    error_str = str(e).lower()
                    if "password" in error_str or "encrypt" in error_str:
                        raise ValueError(f"PDF is password-protected. Please provide the 'password' parameter to extract content.")
                    raise

                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text() or ""

                        # If very little text, try OCR
                        if len(page_text.strip()) < 50 and self._ocr_enabled and self._check_ocr_available():
                            ocr_text = self._ocr_pdf_page_pypdf2(path, page_num)
                            if ocr_text:
                                page_text = ocr_text
                                extraction_method = "ocr"

                        if page_text.strip():
                            text_content.append(page_text)
                    except Exception as e:
                        logger.warning(f"Error reading PDF page {page_num}: {e}")

        except ValueError:
            # Re-raise our custom ValueError messages
            raise
        except Exception as e:
            error_str = str(e).lower()
            if "password" in error_str or "encrypted" in error_str:
                raise ValueError(f"PDF is password-protected. Please provide the 'password' parameter to extract content.")
            logger.error(f"PDF extraction failed: {e}")
            raise ValueError(f"Failed to extract PDF content: {str(e)[:200]}")

        return '\n\n'.join(text_content), extraction_method

    def _ocr_pdf_page(self, page) -> Optional[str]:
        """OCR a single PDF page using pymupdf"""
        try:
            import fitz
            import pytesseract
            from PIL import Image

            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

            text = pytesseract.image_to_string(img, lang=self._ocr_language)
            return text if text.strip() else None
        except Exception as e:
            logger.warning(f"OCR failed for page: {e}")
            return None

    def _ocr_pdf_page_pypdf2(self, path: Path, page_num: int) -> Optional[str]:
        """OCR a PDF page using pdf2image"""
        try:
            import pytesseract
            from pdf2image import convert_from_path

            images = convert_from_path(str(path), first_page=page_num + 1, last_page=page_num + 1, dpi=200)
            if images:
                text = pytesseract.image_to_string(images[0], lang=self._ocr_language)
                return text if text.strip() else None
        except Exception as e:
            logger.warning(f"OCR failed for page {page_num}: {e}")
        return None

    def _extract_docx(self, path: Path) -> str:
        """Extract text from Word document including tables"""
        from docx import Document

        doc = Document(path)
        content = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)

        # Extract tables if enabled
        if self._extract_tables:
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                table_text = self._format_table(table_data)
                if table_text:
                    content.append(f"\n[Table]\n{table_text}")

        # Extract headers and footers
        for section in doc.sections:
            header = section.header
            if header and header.paragraphs:
                header_text = ' '.join(p.text for p in header.paragraphs if p.text.strip())
                if header_text:
                    content.insert(0, f"[Header] {header_text}")

            footer = section.footer
            if footer and footer.paragraphs:
                footer_text = ' '.join(p.text for p in footer.paragraphs if p.text.strip())
                if footer_text:
                    content.append(f"[Footer] {footer_text}")

        return '\n\n'.join(content)

    def _extract_pptx(self, path: Path) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required for PowerPoint files. Install with: pip install python-pptx")

        prs = Presentation(str(path))
        content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                # Text frames
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Tables
                if self._extract_tables and shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    table_text = self._format_table(table_data)
                    if table_text:
                        slide_text.append(f"[Table]\n{table_text}")

            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[Notes] {notes}")

            content.append('\n'.join(slide_text))

        return '\n\n'.join(content)

    def _extract_excel(self, path: Path) -> str:
        """Extract text from Excel spreadsheet"""
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl required for Excel files. Install with: pip install openpyxl")

        wb = load_workbook(str(path), data_only=True)
        content = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_content = [f"[Sheet: {sheet_name}]"]

            rows = list(sheet.iter_rows(values_only=True))
            if rows:
                table_text = self._format_table(rows)
                if table_text:
                    sheet_content.append(table_text)

            content.append('\n'.join(sheet_content))

        return '\n\n'.join(content)

    def _extract_csv(self, path: Path) -> str:
        """Extract text from CSV file"""
        content = []

        # Try different encodings
        for encoding in ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']:
            try:
                with open(path, 'r', encoding=encoding, newline='') as f:
                    # Detect delimiter
                    sample = f.read(4096)
                    f.seek(0)

                    try:
                        dialect = csv.Sniffer().sniff(sample)
                        reader = csv.reader(f, dialect)
                    except csv.Error:
                        reader = csv.reader(f)

                    rows = list(reader)
                    if rows:
                        content.append(self._format_table(rows))
                    break
            except UnicodeDecodeError:
                continue

        return '\n'.join(content)

    def _extract_html(self, path: Path) -> str:
        """Extract text from HTML file"""
        from bs4 import BeautifulSoup

        try:
            html_content = path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            html_content = path.read_text(encoding='latin-1')

        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove unwanted elements
        for element in soup(["script", "style", "nav", "footer", "header", "aside", "noscript"]):
            element.decompose()

        # Extract text
        text = soup.get_text(separator='\n', strip=True)

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        return '\n'.join(line for line in lines if line)

    def _extract_rtf(self, path: Path) -> str:
        """Extract text from RTF file"""
        try:
            # Simple RTF text extraction
            content = path.read_bytes()

            # Basic RTF parsing - extract text between control words
            import re
            text = content.decode('latin-1', errors='ignore')

            # Remove RTF control words
            text = re.sub(r'\\[a-z]+\d*\s?', '', text)
            text = re.sub(r'[{}]', '', text)
            text = re.sub(r'\\\n', '\n', text)
            text = re.sub(r'\\\'([0-9a-fA-F]{2})', lambda m: chr(int(m.group(1), 16)), text)

            return text.strip()
        except Exception as e:
            logger.warning(f"RTF extraction failed: {e}")
            return ""

    def _extract_image_ocr(self, path: Path) -> str:
        """Extract text from image using OCR"""
        if not self._ocr_enabled:
            logger.info(f"OCR disabled, skipping image: {path}")
            return ""

        if not self._check_ocr_available():
            raise ImportError("OCR requires pytesseract and tesseract. Install tesseract-ocr system package.")

        try:
            import pytesseract
            from PIL import Image

            img = Image.open(path)

            # Convert to RGB if necessary
            if img.mode not in ('L', 'RGB'):
                img = img.convert('RGB')

            text = pytesseract.image_to_string(img, lang=self._ocr_language)
            return text.strip()
        except Exception as e:
            logger.error(f"Image OCR failed for {path}: {e}")
            return ""

    def _extract_text(self, path: Path) -> str:
        """Load plain text or markdown file"""
        try:
            return path.read_text(encoding='utf-8')
        except UnicodeDecodeError:
            return path.read_text(encoding='latin-1')

    def _format_table(self, rows: List[List[Any]]) -> str:
        """Format table data as readable text"""
        if not rows:
            return ""

        # Filter out empty rows
        rows = [row for row in rows if any(cell for cell in row if cell is not None and str(cell).strip())]
        if not rows:
            return ""

        # Convert all cells to strings
        str_rows = [[str(cell) if cell is not None else "" for cell in row] for row in rows]

        # Calculate column widths
        if not str_rows:
            return ""

        max_cols = max(len(row) for row in str_rows)
        col_widths = [0] * max_cols

        for row in str_rows:
            for i, cell in enumerate(row):
                if i < max_cols:
                    col_widths[i] = max(col_widths[i], len(cell[:50]))  # Limit cell width

        # Format as simple table
        lines = []
        for row in str_rows:
            padded = [cell[:50].ljust(col_widths[i]) if i < len(row) else "" for i, cell in enumerate(row)]
            lines.append(" | ".join(padded).strip())

        return '\n'.join(lines)
